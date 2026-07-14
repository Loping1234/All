import pptx
import json
import sys

def extract_pptx_text(file_path):
    prs = pptx.Presentation(file_path)
    slides_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_content.append({
            "slide_number": i + 1,
            "content": "\n".join(slide_text)
        })
    return slides_content

if __name__ == "__main__":
    file_path = "Internship_Review_Pt3.pptx"
    try:
        content = extract_pptx_text(file_path)
        print(json.dumps(content, indent=4))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
