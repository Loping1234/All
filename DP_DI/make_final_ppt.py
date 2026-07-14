from pathlib import Path
import math

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TASK_ROOT = ROOT.parent
TEMPLATE = Path(r"C:\Users\PRANAY\Downloads\OneDrive_1_5-8-2026\Final Presentation Format.pptx")
OUT_DIR = ROOT / "output"
ASSET_DIR = ROOT / "ppt_assets"
REPORT_IMAGES = ASSET_DIR / "report_images"
LOGO = ASSET_DIR / "template" / "image1.jpg"
OUTPUT = OUT_DIR / "A149_PranayKumar_Final_Presentation_Revora.pptx"


NAVY = RGBColor(9, 15, 34)
INK = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)
RED = RGBColor(164, 28, 48)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
PURPLE = RGBColor(124, 58, 237)
TEAL = RGBColor(13, 148, 136)
WHITE = RGBColor(255, 255, 255)


def rgb_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def clear_template_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=BORDER, width=1.0, transparency=0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def textbox(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            font="Aptos", valign=MSO_ANCHOR.TOP, margin=0.02, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rich_textbox(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold, run_size in runs:
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(run_size or size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_header(slide, slide_no, title=None):
    slide.shapes.add_picture(str(LOGO), Inches(0.55), Inches(0.22), height=Inches(0.34))
    textbox(slide, "School of Technology Management & Engineering | NMIMS Navi Mumbai",
            1.45, 0.24, 5.9, 0.26, size=8.7, color=MUTED, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if title:
        textbox(slide, title, 7.6, 0.22, 4.65, 0.28, size=8.5, color=MUTED,
                align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.035))
    set_fill(line, RED)
    set_line(line, RED, 0)
    textbox(slide, f"{slide_no:02d}", 12.62, 7.22, 0.38, 0.18, size=7.8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, text, subtitle=None):
    textbox(slide, text, 0.72, 0.78, 8.2, 0.52, size=25, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.34), Inches(0.78), Inches(0.05))
    set_fill(accent, RED)
    set_line(accent, RED, 0)
    if subtitle:
        textbox(slide, subtitle, 0.72, 1.46, 9.4, 0.34, size=11.5, color=SLATE)


def rounded_box(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 0.8)
    return shape


def pill(slide, text, x, y, w, color, bg=None, size=10, bold=True):
    bg = bg or RGBColor(241, 245, 249)
    sh = rounded_box(slide, x, y, w, 0.28, bg, color)
    sh.line.width = Pt(0.6)
    textbox(slide, text, x + 0.06, y + 0.035, w - 0.12, 0.18, size=size, color=color,
            bold=bold, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return sh


def bullet_list(slide, items, x, y, w, h, size=12, color=SLATE, bullet_color=RED, gap=0.34):
    for idx, item in enumerate(items):
        yy = y + idx * gap
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.08), Inches(0.08), Inches(0.08))
        set_fill(dot, bullet_color)
        set_line(dot, bullet_color, 0)
        textbox(slide, item, x + 0.18, yy, w - 0.18, 0.28, size=size, color=color)


def metric(slide, label, value, note, x, y, w, h, color=BLUE):
    rounded_box(slide, x, y, w, h, WHITE, BORDER)
    textbox(slide, label.upper(), x + 0.16, y + 0.14, w - 0.32, 0.18, size=7.5, color=MUTED, bold=True)
    textbox(slide, value, x + 0.16, y + 0.37, w - 0.32, 0.38, size=21, color=color, bold=True)
    textbox(slide, note, x + 0.16, y + 0.82, w - 0.32, 0.38, size=8.5, color=SLATE)


def diagram_node(slide, title, note, x, y, w, h, color=BLUE, fill=WHITE):
    rounded_box(slide, x, y, w, h, fill, BORDER)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(marker, color)
    set_line(marker, color, 0)
    textbox(slide, title, x + 0.18, y + 0.12, w - 0.28, 0.25, size=11, color=INK, bold=True)
    textbox(slide, note, x + 0.18, y + 0.42, w - 0.28, h - 0.48, size=8.5, color=SLATE)


def connector(slide, x1, y1, x2, y2, color=MUTED):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.2)
    return line


def chevron(slide, x, y, color=MUTED, rotate=0):
    ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.22), Inches(0.22))
    set_fill(ch, color)
    set_line(ch, color, 0)
    ch.rotation = rotate
    return ch


def add_picture_fit(slide, image_path, x, y, w, h, border=True):
    from PIL import Image

    path = Path(image_path)
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        height = h
        width = h * img_ratio
        left = x - (width - w) / 2
        top = y
    else:
        width = w
        height = w / img_ratio
        left = x
        top = y - (height - h) / 2
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    if border:
        frame = rounded_box(slide, x, y, w, h, RGBColor(255, 255, 255), BORDER, radius=False)
        frame.fill.background()
        frame.line.width = Pt(1)
    return pic


def add_table(slide, rows, x, y, w, h, header_fill=NAVY, font_size=8.2):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = header_fill if r == 0 else (RGBColor(248, 250, 252) if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size if r else font_size + 0.2)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else (INK if c == 0 else SLATE)
    return table_shape


def section_label(slide, label, x=0.72, y=0.62, color=RED):
    pill(slide, label, x, y, 1.5, color, RGBColor(254, 242, 242), size=8.5)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.02))
    set_fill(band, NAVY)
    set_line(band, NAVY, 0)
    slide.shapes.add_picture(str(LOGO), Inches(0.62), Inches(0.22), height=Inches(0.47))
    textbox(slide, "School of Technology Management & Engineering | NMIMS Navi Mumbai",
            1.78, 0.32, 6.7, 0.28, size=9.5, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    pill(slide, "INDUSTRY PROJECT", 10.74, 0.31, 1.65, RED, RGBColor(255, 242, 242), size=8.5)

    textbox(slide, "Revora", 0.82, 1.55, 4.6, 0.76, size=42, color=INK, bold=True)
    textbox(slide, "Dynamic Pricing & Demand Intelligence System",
            0.84, 2.26, 7.4, 0.5, size=22, color=SLATE, bold=False)
    textbox(slide, "A MERN-based decision-support platform that turns messy sales data into explainable, guarded price recommendations.",
            0.86, 2.95, 7.4, 0.55, size=13.2, color=SLATE)

    metric(slide, "Decision rows", "157,519", "Generated decision-quality artifact", 8.55, 1.55, 1.7, 1.25, BLUE)
    metric(slide, "Accuracy", "86.28%", "Time-based holdout", 10.45, 1.55, 1.7, 1.25, GREEN)
    metric(slide, "Macro F1", "0.7196", "Class-balanced metric", 8.55, 3.03, 1.7, 1.25, PURPLE)
    metric(slide, "Build", "Passed", "Frontend production build", 10.45, 3.03, 1.7, 1.25, AMBER)

    rows = [
        ("Student", "Pranay Kumar | Roll No. A149 | SAP 70022200357"),
        ("Program", "B.Tech Computer Engineering | Semester VIII | AY 2025-26"),
        ("Internal Mentor", "Dr. Preeti Gupta, STME NMIMS"),
        ("Industry Mentor", "Mr. Kanwar Pal Singh Rathore, Tech Biz Solution"),
        ("Company", "Tech Biz Solution, Jodhpur"),
    ]
    y = 4.48
    for label, value in rows:
        textbox(slide, label.upper(), 0.86, y, 1.35, 0.22, size=7.5, color=MUTED, bold=True)
        textbox(slide, value, 2.28, y - 0.02, 6.7, 0.26, size=10.8, color=INK)
        y += 0.38

    # Decision flow motif.
    flow = [("CSV/API", BLUE), ("Quality", TEAL), ("Models", PURPLE), ("Guardrails", AMBER), ("Decision", GREEN)]
    x = 1.0
    for idx, (name, color) in enumerate(flow):
        rounded_box(slide, x, 6.43, 1.35, 0.42, RGBColor(248, 250, 252), color)
        textbox(slide, name, x + 0.1, 6.53, 1.15, 0.17, size=8.6, color=color, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(flow) - 1:
            chevron(slide, x + 1.46, 6.53, MUTED)
        x += 1.72
    return slide


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 2, "Problem Statement")
    add_title(slide, "Problem Statement", "Pricing data exists, but the decision workflow is usually fragmented.")

    textbox(slide, "Core issue", 0.78, 2.04, 1.6, 0.28, size=11, color=RED, bold=True)
    textbox(slide, "Manual price changes can ignore demand sensitivity, inventory risk, competitor gaps, and data quality.",
            0.78, 2.35, 5.25, 0.8, size=19, color=INK, bold=True)
    textbox(slide, "The result is not just a bad number. It is an unaudited business decision with unclear assumptions.",
            0.8, 3.28, 5.0, 0.42, size=12, color=SLATE)

    problems = [
        ("Messy spreadsheet inputs", "duplicate rows, missing cost, inconsistent products", BLUE),
        ("Hidden demand response", "price changes are not tied to quantity movement", PURPLE),
        ("Unsafe recommendations", "no guardrails for below-cost or extreme prices", AMBER),
    ]
    x = 6.62
    for i, (head, note, color) in enumerate(problems):
        diagram_node(slide, head, note, x, 1.84 + i * 1.12, 5.05, 0.82, color)
        if i < len(problems) - 1:
            connector(slide, x + 2.55, 2.66 + i * 1.12, x + 2.55, 2.9 + i * 1.12)
            chevron(slide, x + 2.45, 2.7 + i * 1.12, MUTED, rotate=90)

    rounded_box(slide, 0.78, 5.28, 11.0, 1.02, RGBColor(255, 251, 235), RGBColor(252, 211, 77))
    textbox(slide, "Why it matters", 1.02, 5.45, 1.7, 0.22, size=10.2, color=AMBER, bold=True)
    bullet_list(slide, [
        "Retail and service businesses need price decisions that protect margin and inventory.",
        "Managers must understand the reason for each recommendation before trusting it.",
        "A final-year implementation project must show real workflows, not only a model."
    ], 1.05, 5.66, 10.0, 0.5, size=8.9, bullet_color=AMBER, gap=0.18)
    return slide


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 3, "Objectives")
    add_title(slide, "Objectives & Expected Outcomes", "Six practical goals guide the Revora implementation.")

    goals = [
        ("01", "Authenticated pricing workspace", "Admin/demo login, JWT sessions, workspace-aware actions.", BLUE),
        ("02", "CSV ingestion and quality gate", "Mapping, staging, duplicate checks, row issues, data fitness.", TEAL),
        ("03", "Demand intelligence", "Price-response models using product, segment, inventory, competitor context.", PURPLE),
        ("04", "Scenario simulation", "Demand, revenue, profit, confidence, and warnings for tested prices.", AMBER),
        ("05", "Guarded recommendations", "Good range, avoid range, competitor distance, business objective.", GREEN),
        ("06", "Results export", "Dashboards, ML Decision Space, XLSX/CSV reports for review.", RED),
    ]
    for i, (num, title, note, color) in enumerate(goals):
        col = i % 2
        row = i // 2
        x = 0.88 + col * 6.05
        y = 1.92 + row * 1.38
        rounded_box(slide, x, y, 5.55, 0.98, WHITE, BORDER)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.18), Inches(0.5), Inches(0.5))
        set_fill(circ, color)
        set_line(circ, color, 0)
        textbox(slide, num, x + 0.26, y + 0.32, 0.38, 0.12, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, title, x + 0.86, y + 0.16, 4.35, 0.24, size=12.3, color=INK, bold=True)
        textbox(slide, note, x + 0.86, y + 0.48, 4.38, 0.28, size=9.2, color=SLATE)

    rounded_box(slide, 1.02, 6.18, 11.05, 0.48, RGBColor(240, 253, 244), RGBColor(187, 247, 208))
    textbox(slide, "Expected outcome: an examiner-ready working product with visible UI, data pipeline, explainable recommendations, measurable ML results, and exportable reports.",
            1.22, 6.32, 10.65, 0.17, size=10.2, color=RGBColor(22, 101, 52), bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 4, "Industry Relevance")
    add_title(slide, "Industry Relevance / Motivation", "Revora is positioned as decision support for real business pricing, not blind automation.")

    x0 = 0.92
    phases = [
        ("Foundation", "Node, Express, MongoDB, auth", BLUE),
        ("TaskOps", "RBAC, task boards, real-time workflows", TEAL),
        ("NewsTv19", "RSS ingestion, search, production codebase", AMBER),
        ("Revora", "pricing intelligence and decision quality", RED),
    ]
    for i, (title, note, color) in enumerate(phases):
        x = x0 + i * 3.05
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.08), Inches(0.48), Inches(0.48))
        set_fill(circ, color)
        set_line(circ, color, 0)
        textbox(slide, str(i + 1), x + 0.16, 2.22, 0.16, 0.1, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, title, x - 0.18, 2.72, 1.1, 0.24, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, note, x - 0.48, 3.05, 1.68, 0.42, size=8.2, color=SLATE, align=PP_ALIGN.CENTER)
        if i < 3:
            connector(slide, x + 0.52, 2.32, x + 2.85, 2.32)
            chevron(slide, x + 2.7, 2.2, MUTED)

    rows = [
        ["Business need", "Manual pricing needs margin, demand, inventory, and competitor context."],
        ["Practical applicability", "Retail catalog managers can import CSVs and review safe ranges before acting."],
        ["Market fit", "Small teams want BI-style guidance without requiring a data-science operator."],
        ["Project motivation", "Convert internship backend/frontend learning into one analytics product."],
    ]
    add_table(slide, [["Angle", "How Revora answers it"]] + rows, 0.95, 4.08, 11.2, 1.65, font_size=8.7)
    textbox(slide, "Panel answer cue", 0.98, 6.12, 1.5, 0.18, size=8.5, color=RED, bold=True)
    textbox(slide, "The system keeps a human in control: it recommends, explains, warns, and exports evidence.",
            2.32, 6.08, 8.9, 0.25, size=12.5, color=INK, bold=True)
    return slide


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 5, "Existing System")
    add_title(slide, "Literature Survey / Existing System", "The gap is not dynamic pricing as a concept; it is trusted execution for messy business data.")

    rows = [
        ["Approach", "Strength", "Limitation", "Revora response"],
        ["Manual Excel pricing", "Fast and familiar", "No model-readiness or guardrail trail", "CSV import + quality labels"],
        ["Generic BI dashboards", "Good visualization", "Shows what happened, not what price to test", "Simulation + recommendation engine"],
        ["Black-box price tools", "Automated suggestions", "Weak explainability for academic/SME review", "Calculation steps + warnings"],
        ["Research models", "Strong theory", "Often assume cleaner data than real SMEs have", "Staged imports + safe ranges"],
    ]
    add_table(slide, rows, 0.78, 1.92, 11.78, 2.38, header_fill=NAVY, font_size=7.6)

    rounded_box(slide, 0.88, 4.78, 5.5, 1.2, RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    textbox(slide, "Literature motivation", 1.08, 4.98, 2.1, 0.22, size=10.5, color=BLUE, bold=True)
    bullet_list(slide, [
        "Dynamic pricing learns from historical demand and price experimentation.",
        "Limited experimentation increases the need for confidence and safeguards."
    ], 1.08, 5.3, 4.9, 0.5, size=9.5, bullet_color=BLUE, gap=0.25)

    rounded_box(slide, 6.75, 4.78, 5.5, 1.2, RGBColor(240, 253, 244), RGBColor(187, 247, 208))
    textbox(slide, "Design implication", 6.95, 4.98, 2.0, 0.22, size=10.5, color=GREEN, bold=True)
    bullet_list(slide, [
        "Treat the model as decision support.",
        "Expose assumptions, constraints, and exportable evidence."
    ], 6.95, 5.3, 4.85, 0.5, size=9.5, bullet_color=GREEN, gap=0.25)
    return slide


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 6, "Proposed System")
    add_title(slide, "Proposed System Architecture", "A layered MERN workflow from import to auditable pricing decision.")

    # Left-to-right architecture diagram.
    diagram_node(slide, "Users", "admin / analyst / examiner demo", 0.78, 3.0, 1.5, 0.75, RED)
    diagram_node(slide, "React + Vite UI", "workspaces, dashboards, upload, simulator", 2.72, 2.1, 2.05, 1.0, BLUE)
    diagram_node(slide, "Express API", "auth, routes, validation, reports", 5.25, 2.1, 2.0, 1.0, TEAL)
    diagram_node(slide, "MongoDB", "products, sales, batches, models, recommendations", 7.78, 2.1, 2.05, 1.0, PURPLE)
    diagram_node(slide, "Python ML Space", "offline decision-quality training and prediction", 10.32, 2.1, 2.15, 1.0, AMBER)

    diagram_node(slide, "CSV/API ingestion", "mapping, staging, duplicate checks", 2.72, 4.05, 2.05, 0.9, BLUE)
    diagram_node(slide, "Quality gate", "fitness score, readiness, row issues", 5.25, 4.05, 2.0, 0.9, GREEN)
    diagram_node(slide, "Pricing engine", "demand model, simulation, good/avoid ranges", 7.78, 4.05, 2.05, 0.9, RED)
    diagram_node(slide, "Reports", "XLSX/CSV exports and screenshot evidence", 10.32, 4.05, 2.15, 0.9, TEAL)

    for x1, y1, x2, y2 in [
        (2.28, 3.38, 2.72, 2.58), (4.77, 2.58, 5.25, 2.58), (7.25, 2.58, 7.78, 2.58),
        (9.83, 2.58, 10.32, 2.58), (3.75, 3.1, 3.75, 4.05), (4.77, 4.5, 5.25, 4.5),
        (7.25, 4.5, 7.78, 4.5), (9.83, 4.5, 10.32, 4.5)
    ]:
        connector(slide, x1, y1, x2, y2)

    rounded_box(slide, 0.9, 5.72, 11.34, 0.52, RGBColor(248, 250, 252), BORDER)
    textbox(slide, "Key flow: Import data -> verify fitness -> fit/compare models -> simulate price candidates -> apply trust policy -> store recommendation -> export report.",
            1.08, 5.89, 10.98, 0.17, size=9.8, color=INK, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 7, "Technology Stack")
    add_title(slide, "Technology Stack", "Chosen for a full-stack product workflow, not only a model demo.")

    stacks = [
        ("Frontend", ["React 19", "Vite 6", "Tailwind CSS 4", "Lucide icons", "Custom chart components"], BLUE),
        ("Backend/API", ["Node.js", "Express", "Mongoose", "Multer", "csv-parse", "ExcelJS"], TEAL),
        ("Data & ML", ["MongoDB", "Python ML pipeline", "scikit-learn RandomForestClassifier", "Generated feature dataset"], PURPLE),
        ("Verification", ["Unit reliability tests", "ML self-test", "Production frontend build", "Browser/API evidence"], GREEN),
    ]
    for i, (title, items, color) in enumerate(stacks):
        x = 0.78 + (i % 2) * 6.1
        y = 1.92 + (i // 2) * 2.12
        rounded_box(slide, x, y, 5.58, 1.62, WHITE, BORDER)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.58), Inches(0.08)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = color
        slide.shapes[-1].line.color.rgb = color
        textbox(slide, title, x + 0.22, y + 0.22, 2.0, 0.28, size=14, color=INK, bold=True)
        bullet_list(slide, items, x + 0.26, y + 0.62, 5.0, 0.8, size=8.8, bullet_color=color, gap=0.22)

    add_table(slide, [
        ["Why this stack?", "Reason"],
        ["MERN", "Single JavaScript product layer for UI, API, and data models."],
        ["MongoDB", "Flexible document storage for variable sales/import records."],
        ["Python ML", "Keeps offline model training separate from explainable pricing API."],
    ], 1.0, 6.0, 10.95, 0.75, header_fill=NAVY, font_size=7.5)
    return slide


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 8, "Methodology")
    add_title(slide, "Methodology: Data-to-Decision Workflow", "The implementation follows an auditable path from raw records to final recommendation.")

    steps = [
        ("1", "Import", "CSV/API rows", BLUE),
        ("2", "Map", "headers + products", TEAL),
        ("3", "Check", "quality + readiness", GREEN),
        ("4", "Model", "price response", PURPLE),
        ("5", "Simulate", "demand/revenue/profit", AMBER),
        ("6", "Recommend", "guardrails + export", RED),
    ]
    x = 0.82
    for i, (num, title, note, color) in enumerate(steps):
        rounded_box(slide, x, 2.35, 1.56, 1.08, WHITE, color)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.5), Inches(2.04), Inches(0.52), Inches(0.52))
        set_fill(circ, color)
        set_line(circ, WHITE, 1)
        textbox(slide, num, x + 0.68, 2.19, 0.16, 0.1, size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, title, x + 0.15, 2.62, 1.26, 0.2, size=11.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, note, x + 0.12, 2.96, 1.32, 0.28, size=8.3, color=SLATE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            chevron(slide, x + 1.68, 2.78, MUTED)
        x += 2.02

    lanes = [
        ("Quality controls", "duplicate detection, missing-cost checks, stockout exclusion, grouped demand points"),
        ("Model controls", "linear/log-log/context-aware choices, backtest summary, reliability labels"),
        ("Decision controls", "below-cost blocking, extreme-price warnings, good range, avoid range, competitor distance"),
    ]
    for i, (head, note) in enumerate(lanes):
        y = 4.55 + i * 0.55
        rounded_box(slide, 1.05, y, 11.05, 0.34, RGBColor(248, 250, 252), BORDER)
        textbox(slide, head, 1.22, y + 0.09, 1.55, 0.12, size=8.5, color=INK, bold=True)
        textbox(slide, note, 2.92, y + 0.09, 8.75, 0.12, size=8.5, color=SLATE)
    return slide


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 9, "Implementation")
    add_title(slide, "Implementation Modules", "The product is split into reviewable modules with clear data ownership.")

    modules = [
        ("Import Batch", "tracks source, status, mapping, row counts", BLUE),
        ("Sales Data", "normalized price, quantity, segment, date", TEAL),
        ("Product", "SKU, category, cost, inventory, readiness", GREEN),
        ("Demand Model", "model type, coefficients, backtest metrics", PURPLE),
        ("Recommendation", "tested prices, expected outcomes, trust label", RED),
        ("Assistant Decision", "chat-derived pricing decisions and labels", AMBER),
    ]
    coords = [(0.85, 2.0), (3.05, 2.0), (5.25, 2.0), (2.0, 4.1), (4.2, 4.1), (6.4, 4.1)]
    for (title, note, color), (x, y) in zip(modules, coords):
        diagram_node(slide, title, note, x, y, 1.86, 0.74, color)

    # Relationship hints.
    for x1, y1, x2, y2 in [
        (1.78, 2.74, 3.05, 2.37), (4.91, 2.37, 5.25, 2.37),
        (5.93, 2.74, 5.13, 4.1), (3.92, 2.74, 2.93, 4.1),
        (3.86, 4.47, 4.2, 4.47), (6.06, 4.47, 6.4, 4.47)
    ]:
        connector(slide, x1, y1, x2, y2)

    rounded_box(slide, 8.75, 1.98, 3.42, 3.0, RGBColor(248, 250, 252), BORDER)
    textbox(slide, "Real functionality implemented", 8.98, 2.22, 2.75, 0.24, size=13.2, color=INK, bold=True)
    bullet_list(slide, [
        "CSV upload with field mapping and staged commit",
        "Workspace reset and active dataset controls",
        "Pricing insights and model-readiness checks",
        "Scenario planner and objective-based recommendation",
        "ML Decision Space with trained classifier metrics",
        "XLSX + CSV report exports"
    ], 9.02, 2.65, 2.85, 1.6, size=8.8, bullet_color=RED, gap=0.3)

    textbox(slide, "Most important viva point", 0.95, 6.02, 2.1, 0.18, size=8.4, color=RED, bold=True)
    textbox(slide, "The recommendation is not stored as a number only; it stores evidence, assumptions, warnings, and tested alternatives.",
            2.75, 5.98, 8.88, 0.22, size=10.5, color=INK, bold=True)
    return slide


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 10, "Implementation Evidence")
    add_title(slide, "Implementation Screenshots: Workspace & Import", "Real UI evidence from the Revora local browser session.")

    add_picture_fit(slide, REPORT_IMAGES / "report_13.jpg", 0.78, 1.76, 5.92, 3.32)
    add_picture_fit(slide, REPORT_IMAGES / "report_14.jpg", 6.94, 1.76, 5.92, 3.32)
    pill(slide, "Home workspace", 1.0, 5.25, 1.45, BLUE, RGBColor(239, 246, 255), size=8.5)
    textbox(slide, "Revenue snapshot, model-ready products, active dataset, and dashboard navigation.",
            2.58, 5.28, 3.78, 0.16, size=8.8, color=SLATE)
    pill(slide, "Sales data upload", 7.18, 5.25, 1.65, TEAL, RGBColor(240, 253, 250), size=8.5)
    textbox(slide, "CSV selection, preview mapping, quality review, and commit pipeline.",
            8.96, 5.28, 3.38, 0.16, size=8.8, color=SLATE)
    rounded_box(slide, 1.05, 6.06, 11.0, 0.48, RGBColor(248, 250, 252), BORDER)
    textbox(slide, "Demo talking point: start with dataset reset, upload sales_diverse.csv, then show how Revora converts rows into pricing-ready product groups.",
            1.24, 6.21, 10.62, 0.16, size=9.3, color=INK, bold=True, align=PP_ALIGN.CENTER)
    return slide


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 11, "Feature Evidence")
    add_title(slide, "Implementation Screenshots: Insights, ML & Reports", "Three examiner-visible surfaces prove the pipeline is not just backend code.")

    add_picture_fit(slide, REPORT_IMAGES / "report_15.jpg", 0.78, 1.74, 3.75, 2.12)
    add_picture_fit(slide, REPORT_IMAGES / "report_16.jpg", 4.78, 1.74, 3.75, 2.12)
    add_picture_fit(slide, REPORT_IMAGES / "report_17.jpg", 8.78, 1.74, 3.75, 2.12)
    features = [
        ("Pricing Insights", "30 ready combinations, product selection, price-response creation", BLUE),
        ("ML Decision Space", "157,519 decision rows, 86.3% UI accuracy, feature signals", PURPLE),
        ("Reports & Export", "XLSX/CSV reports for dashboard, products, sales, insights", GREEN),
    ]
    for i, (title, note, color) in enumerate(features):
        x = 0.86 + i * 4.0
        rounded_box(slide, x, 4.28, 3.5, 1.1, WHITE, BORDER)
        textbox(slide, title, x + 0.18, 4.48, 2.9, 0.22, size=12.5, color=color, bold=True)
        textbox(slide, note, x + 0.18, 4.82, 3.0, 0.36, size=8.8, color=SLATE)
    textbox(slide, "Real functionality over theory", 0.98, 6.05, 2.3, 0.18, size=8.5, color=RED, bold=True)
    textbox(slide, "The screenshots align directly with the demo path: create insight -> inspect ML evidence -> export report.",
            3.0, 6.0, 8.6, 0.25, size=11.2, color=INK, bold=True)
    return slide


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 12, "Results")
    add_title(slide, "Results & Analysis", "The final report includes measured ML output and local verification evidence.")

    metric(slide, "Total rows", "157,519", "Generated decision-quality dataset", 0.78, 1.72, 2.0, 1.05, BLUE)
    metric(slide, "Training rows", "126,015", "Used for model fitting", 3.02, 1.72, 2.0, 1.05, TEAL)
    metric(slide, "Testing rows", "31,504", "Time-based holdout", 5.26, 1.72, 2.0, 1.05, AMBER)
    metric(slide, "Accuracy", "86.28%", "Overall classification", 7.5, 1.72, 2.0, 1.05, GREEN)
    metric(slide, "Macro F1", "0.7196", "Balanced class quality", 9.74, 1.72, 2.0, 1.05, PURPLE)

    chart_data = CategoryChartData()
    chart_data.categories = ["Terrible", "Bad", "Neutral", "Good", "Terrific"]
    chart_data.add_series("Rows", (36619, 76215, 5373, 3058, 36254))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.88), Inches(3.28), Inches(5.55), Inches(2.52), chart_data).chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.font.size = Pt(7)
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE
    textbox(slide, "Decision-quality label distribution", 0.95, 3.05, 3.3, 0.22, size=10.2, color=INK, bold=True)

    rows = [
        ["Check", "Method", "Result"],
        ["ML pipeline", "npm run test:ml-pipeline", "Passed"],
        ["Reliability suite", "npm run test:reliability", "Passed"],
        ["Frontend build", "npm run build", "Passed with large chunk warning"],
        ["UI evidence", "Browser screenshots", "Captured"],
    ]
    add_table(slide, rows, 6.88, 3.28, 5.45, 2.52, header_fill=NAVY, font_size=7.5)
    textbox(slide, "Interpretation", 6.94, 5.98, 1.0, 0.18, size=8.6, color=RED, bold=True)
    textbox(slide, "The ML result is prototype evidence: useful for viva demonstration, but real client transaction data is required before autonomous production pricing.",
            7.86, 5.94, 4.3, 0.26, size=8.9, color=SLATE)
    return slide


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 13, "Challenges")
    add_title(slide, "Challenges Faced, Risk Controls & Failure Testing", "The hardest part was making recommendations safe enough to explain.")

    rows = [
        ["Challenge", "What was done", "Panel answer"],
        ["Varied CSV data", "field mapping, staging, duplicate checks", "Bad data is stopped before modeling"],
        ["Sparse price variation", "readiness labels and grouped demand thresholds", "Some products remain summary-only"],
        ["Unsafe prices", "cost guardrails, range warnings, avoid band", "Below-cost and extreme prices are flagged"],
        ["API/server failure", "health route, user-visible errors, local demo backup", "Demo can continue with screenshots/recording"],
        ["Performance growth", "build check plus future code splitting", "Large chunk warning is known and scoped"],
    ]
    add_table(slide, rows, 0.78, 1.92, 11.78, 2.7, header_fill=NAVY, font_size=7.4)

    controls = [
        ("Security", "JWT auth, protected routes, API-key ingestion path", BLUE),
        ("Scalability", "MongoDB indexes, staged imports, report exports", GREEN),
        ("Trust", "fitness score, warnings, model limits, human review", RED),
    ]
    for i, (title, note, color) in enumerate(controls):
        x = 1.02 + i * 3.95
        rounded_box(slide, x, 5.18, 3.35, 0.86, RGBColor(248, 250, 252), color)
        textbox(slide, title, x + 0.16, 5.38, 1.1, 0.18, size=11.2, color=color, bold=True)
        textbox(slide, note, x + 0.16, 5.66, 2.8, 0.2, size=8.5, color=SLATE)
    return slide


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 14, "Future Scope")
    add_title(slide, "Scalability & Future Scope", "Next work moves Revora from examiner-ready prototype toward production decision support.")

    roadmap = [
        ("Now", "Local MERN demo, CSV workflow, guarded recommendations", BLUE),
        ("Near term", "Real client transaction training + A/B testing", GREEN),
        ("Scale", "ERP/POS imports, marketplace connectors, scheduled reports", TEAL),
        ("Production", "cloud deployment, stronger policies, code splitting", RED),
    ]
    for i, (phase, note, color) in enumerate(roadmap):
        x = 0.9 + i * 3.0
        rounded_box(slide, x, 2.18, 2.46, 1.12, WHITE, color)
        textbox(slide, phase, x + 0.2, 2.43, 2.05, 0.24, size=14.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, note, x + 0.18, 2.78, 2.1, 0.32, size=8.6, color=SLATE, align=PP_ALIGN.CENTER)
        if i < 3:
            chevron(slide, x + 2.6, 2.62, MUTED)

    rows = [
        ["Priority", "Enhancement", "Benefit"],
        ["High", "Train with real client data", "Improves reliability and production relevance"],
        ["High", "Controlled price experiments", "Separates correlation from elasticity"],
        ["High", "Route-level code splitting", "Improves frontend load performance"],
        ["Medium", "ERP/POS imports", "Reduces manual CSV dependency"],
        ["Medium", "Fairness and max-change policies", "Prevents unsafe automation"],
    ]
    add_table(slide, rows, 1.0, 4.06, 11.1, 1.75, header_fill=NAVY, font_size=7.5)
    textbox(slide, "Commercial opportunity", 1.02, 6.25, 1.7, 0.18, size=8.5, color=RED, bold=True)
    textbox(slide, "Revora can become a lightweight pricing cockpit for SMEs that need evidence-backed decisions before they can invest in enterprise pricing suites.",
            2.7, 6.21, 8.86, 0.22, size=10.1, color=INK, bold=True)
    return slide


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 15, "Conclusion")
    add_title(slide, "Conclusion, Demo Plan & References", "Final outcome: an auditable pricing-intelligence product with real implementation evidence.")

    rounded_box(slide, 0.86, 1.88, 5.55, 1.4, RGBColor(240, 253, 244), RGBColor(187, 247, 208))
    textbox(slide, "Conclusion", 1.08, 2.08, 1.4, 0.25, size=13.5, color=GREEN, bold=True)
    textbox(slide, "Revora demonstrates a complete data-to-price workflow: import, clean, model, simulate, recommend, explain, and export.",
            1.08, 2.48, 4.92, 0.44, size=13.2, color=INK, bold=True)

    rounded_box(slide, 6.84, 1.88, 5.55, 1.4, RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    textbox(slide, "Demo backup", 7.06, 2.08, 1.5, 0.25, size=13.5, color=BLUE, bold=True)
    bullet_list(slide, [
        "Working local build and browser path",
        "Offline screenshots in this deck",
        "Recorded demo backup recommended",
        "Internet backup for deployment/docs"
    ], 7.1, 2.43, 4.6, 0.72, size=9.2, bullet_color=BLUE, gap=0.23)

    rows = [
        ["Reference", "Source"],
        ["[1]", "Node.js, Express, React, Vite, MongoDB, Mongoose documentation"],
        ["[2]", "Tailwind CSS and scikit-learn documentation"],
        ["[3]", "den Boer, Dynamic pricing and learning, 2015"],
        ["[4]", "Cheung, Simchi-Levi, Wang, Dynamic pricing with limited experimentation"],
        ["[5]", "Revora source code and final project report, May 2026"],
    ]
    add_table(slide, rows, 0.9, 4.0, 11.45, 1.45, header_fill=NAVY, font_size=7.4)
    textbox(slide, "Viva close", 1.0, 6.08, 1.0, 0.18, size=8.6, color=RED, bold=True)
    textbox(slide, "My contribution: MERN implementation, CSV/data quality workflow, pricing logic, ML Decision Space integration, reports, testing, and documentation.",
            1.9, 6.04, 9.9, 0.22, size=10.3, color=INK, bold=True)
    return slide


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    template_prs = Presentation(str(TEMPLATE))
    prs = Presentation()
    prs.slide_width = template_prs.slide_width
    prs.slide_height = template_prs.slide_height

    for builder in [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
        slide_15,
    ]:
        builder(prs)

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
